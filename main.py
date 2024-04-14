import eel
from googletrans import Translator, LANGUAGES
from docx import Document
import openpyxl
from pptx import Presentation
from plyer import notification
import operator
import sys
import time
import pyautogui
import pyjokes
import pyttsx3
import speech_recognition as sr
import datetime
import os
import cv2
import webbrowser
from requests import get
import wikipedia
import pywhatkit as kit
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import requests
from email.mime.base import MIMEBase
import email
import geocoder
from bs4 import BeautifulSoup
import psutil
import speedtest
import subprocess
import platform
import gamess

# Initialize Eel
eel.init('views')

# Create an instance of the pyttsx3 engine
engine = pyttsx3.init('sapi5')
voices = engine.getProperty('voices')
engine.setProperty('voice', voices[1].id)


# Function to speak and log messages
def speak(message):
    engine.say(message)
    engine.runAndWait()


# Function to log user messages on the JavaScript side
@eel.expose
def log_user_message(message):
    eel.updateUserLog(message)


# Function to log system messages on the JavaScript side
@eel.expose
def log_system_message(message):
    eel.updateSystemLog(message)


# Function to speak and log messages
def speak_and_log(message, message_type="user"):
    if message_type == "user":
        log_user_message(message)
    else:
        log_system_message(message)

    speak(message)


# voice to text
def takecommand():
    r = sr.Recognizer()

    with sr.Microphone() as source:
        speak("Listening....")
        r.pause_threshold = 1
        audio = r.listen(source, timeout=10, phrase_time_limit=5)

    start_time = time.time()  # Record the start time

    try:
        speak("Recognizing...")
        query = r.recognize_google(audio, language='en-in')
        log_system_message(f"user said: {query}")

    except sr.UnknownValueError:
        speak_and_log("Sorry, I didn't catch that.", message_type="system")
        return "none"

    except sr.RequestError as e:
        speak_and_log(f"Could not request results from Google Speech Recognition service; {e}", message_type="system")
        return "none"

    except Exception as e:
        speak_and_log(f"An error occurred: {e}", message_type="system")
        return "none"

    elapsed_time = time.time() - start_time  # Calculate the elapsed time

    if elapsed_time > 60:  # If no speech is detected for more than 60 seconds, raise an exception
        raise TimeoutError("No speech detected for 1 minute")

    return query


# to wish
def wish():
    hour = int(datetime.datetime.now().hour)
    if 0 <= hour < 12:
        speak_and_log("Good Morning")
    elif 12 <= hour < 18:
        speak_and_log("Good Afternoon")
    else:
        speak_and_log("Good Evening")
    speak_and_log("I am Luna. Please tell me how can I help you")


def get_current_location():
    # Get the user's current location based on IP address
    location = geocoder.ip('me')
    return location


# take screenshot
def take_screenshot_and_save():
    # Take a screenshot
    screenshot = pyautogui.screenshot()

    # Define the file name using the current date and time
    timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    file_name = f"screenshot_{timestamp}.png"

    # Save the screenshot to the desktop
    desktop_path = "C:\\Users\\Mayuri\\Pictures\\Screenshots"
    file_path = f"{desktop_path}/{file_name}"

    screenshot.save(file_path)
    speak_and_log(f"Screenshot saved to: {file_path}")


def create_and_open_word_document(file_name):
    document = Document()

    # Add content to the document (modify as needed)
    document.add_heading('New Word Document', level=1)
    document.add_paragraph("This is a new Word document created using Python.")

    # Save the document with the specified file name (without extension)
    file_path = f"D:\\{file_name}.docx"
    document.save(file_path)
    print(f"Word document created successfully: {file_path}")
    speak_and_log(f"Word document created successfully: {file_path}")
    # Open the newly created Word document
    open_word_file(file_path)


def open_word_file(file_path):
    system_platform = platform.system()

    if system_platform == "Windows":
        try:
            subprocess.Popen(["start", "", file_path], shell=True)
            print(f"Opening Word file: {file_path}")
            speak_and_log(f"Opening Word file: {file_path}")

        except Exception as e:
            speak_and_log(f"Error opening Word file: {e}")
    else:
        speak_and_log("Opening Word files not supported on this platform.")


def open_word():
    system_platform = platform.system()

    if system_platform == "Windows":
        try:
            subprocess.Popen(["start", "winword"], shell=True)
            speak_and_log("Opening Microsoft Word...")
        except Exception as e:
            print(f"Error opening Microsoft Word: {e}")
            speak_and_log(f"Error opening Microsoft Word: {e}")

    else:
        speak_and_log("Microsoft Word not supported on this platform.")


def open_excel():
    system_platform = platform.system()

    if system_platform == "Windows":
        try:
            subprocess.Popen(["start", "excel"], shell=True)
            speak_and_log("Opening Microsoft Excel...")
        except Exception as e:
            speak_and_log(f"Error opening Microsoft Excel: {e}")
    else:
        speak_and_log("Microsoft Excel not supported on this platform.")


def create_and_open_excel(file_name):
    # Create a new Excel workbook
    workbook = openpyxl.Workbook()

    # Get the active sheet
    sheet = workbook.active

    # Add some data to the sheet
    sheet['A1'] = 'Hello'
    sheet['B1'] = 'World'

    # Save the workbook
    file_path = f"D:\\{file_name}.xlsx"
    workbook.save(file_path)
    print(f"Excel file created successfully: {file_path}")
    speak_and_log(f"Excel file created successfully: {file_path}")

    # Open the newly created Excel file
    open_excel_file(file_path)


def open_excel_file(file_path):
    system_platform = platform.system()

    if system_platform == "Windows":
        try:
            subprocess.Popen(["start", "", file_path], shell=True)
            speak_and_log(f"Opening Excel file: {file_path}")
        except Exception as e:
            speak_and_log(f"Error opening Excel file: {e}")
    else:
        speak_and_log("Opening Excel files not supported on this platform.")


def open_powerpoint():
    system_platform = platform.system()

    if system_platform == "Windows":
        try:
            subprocess.Popen(["start", "powerpnt"], shell=True)
            speak_and_log("Opening Microsoft PowerPoint...")
        except Exception as e:
            speak_and_log(f"Error opening Microsoft PowerPoint: {e}")
    else:
        speak_and_log("Microsoft PowerPoint not supported on this platform.")


def create_and_open_powerpoint(file_name):
    # Create a new PowerPoint presentation
    presentation = Presentation()

    # Add a slide to the presentation
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title = slide.shapes.title
    title.text = 'Hello, PowerPoint!'

    # Save the presentation
    file_path = f"D:\\{file_name}.pptx"
    presentation.save(file_path)
    print(f"PowerPoint file created successfully: {file_path}")
    speak_and_log(f"PowerPoint file created successfully: {file_path}")

    # Open the newly created PowerPoint file
    open_powerpoint_file(file_path)


def open_powerpoint_file(file_path):
    system_platform = platform.system()

    if system_platform == "Windows":
        try:
            subprocess.Popen(["start", "", file_path], shell=True)
            speak_and_log(f"Opening PowerPoint file: {file_path}")
        except Exception as e:
            speak_and_log(f"Error opening PowerPoint file: {e}")
    else:
        speak_and_log("Opening PowerPoint files not supported on this platform.")


def get_language_code(language_name):
    for code, name in LANGUAGES.items():
        if name.lower() == language_name.lower():
            return code
    return None


def translate_text(text, target_language):
    translator = Translator()
    target_language_code = get_language_code(target_language)

    if target_language_code:
        translation = translator.translate(text, dest=target_language_code)
        return translation.text
    else:
        return "Language not found or supported."


def open_calculator():
    system_platform = platform.system()

    if system_platform == "Windows":
        try:
            subprocess.Popen(["start", "calc"], shell=True)
            print("Opening Calculator...")
        except Exception as e:
            print(f"Error opening Calculator: {e}")
    else:
        print("Calculator not supported on this platform.")


def open_vscode():
    system_platform = platform.system()

    if system_platform == "Windows":
        try:
            subprocess.Popen(["code"], shell=True)
            speak_and_log("Opening Visual Studio Code...")
        except Exception as e:
            speak_and_log(f"Error opening Visual Studio Code: {e}")
    else:
        speak_and_log("Visual Studio Code not supported on this platform.")


def create_and_open_python_file(file_name):
    # Add the '.py' extension to the provided filename
    file_name_with_extension = f"D:\\{file_name}.py"

    # Create a new Python file
    with open(file_name_with_extension, 'w') as file:
        file.write("# Your Python code goes here")

    # Open the file in Python IDLE
    subprocess.Popen(["python", "-m", "idlelib.idle", file_name_with_extension], shell=True)
    print(f"Created and opened Python file '{file_name_with_extension}' in IDLE.")


def find_file(file_name):
    for drive in range(ord('A'), ord('Z') + 1):
        drive_letter = chr(drive) + ":"
        try:
            for root, dirs, files in os.walk(drive_letter):
                if file_name in files:
                    return os.path.join(root, file_name)
        except PermissionError:
            # Handle PermissionError, as some drives may be inaccessible
            pass

    return None


def open_pdf(pdf_file_path):
    if pdf_file_path is not None:
        # Open the PDF file using the default PDF viewer
        webbrowser.open(pdf_file_path)
        print(f"Opening PDF file: {pdf_file_path}")
    else:
        print("Error: PDF file not found in the system.")


# to send email
def sendEmail(recipient_email, subject, body):
    sender_email = "manjarekarmayu25@gmail.com"
    app_password = "ezqb pkfk cuar aayh"

    # Get recipient email address from user input

    message = MIMEMultipart()
    message['From'] = sender_email
    message['To'] = recipient_email
    message['Subject'] = subject
    message.attach(MIMEText(body, 'plain'))

    # Connect to the SMTP server (in this case, Gmail's SMTP server)
    smtp_server = "smtp.gmail.com"
    smtp_port = 587
    server = smtplib.SMTP(smtp_server, smtp_port)
    server.starttls()

    # Log in to your email account using the App Password
    server.login(sender_email, app_password)

    # Send the email
    server.sendmail(sender_email, recipient_email, message.as_string())

    # Quit the server
    server.quit()


# to send file via email
def send_fileemail(receiver_email, file_name):
    sender_email = "manjarekarmayu25@gmail.com"  # Replace with your email
    app_key = "ezqb pkfk cuar aayh"  # Replace with your app-specific key

    file_path = find_file(file_name)

    if file_path is None:
        speak_and_log(f"File '{file_name}' not found.")
        return

    # Set up the MIME
    message = MIMEMultipart()
    message["From"] = sender_email
    message["To"] = receiver_email
    message["Subject"] = "File Sharing"

    # Attach the body to the email
    message.attach(MIMEText("Please find the attached file.", "plain"))

    # Attach the file to the email
    attachment = open(file_path, "rb")
    base = MIMEBase("application", "octet-stream")
    base.set_payload((attachment).read())
    email.encoders.encode_base64(base)
    base.add_header("Content-Disposition", f"attachment; filename={file_name}")
    message.attach(base)

    # Connect to the SMTP server
    smtp_server = "smtp.gmail.com"  # Update this based on your email provider
    smtp_port = 587  # Update this based on your email provider

    with smtplib.SMTP(smtp_server, smtp_port) as server:
        server.starttls()

        # Log in to your email account using app-specific key
        server.login(sender_email, app_key)

        # Send the email
        server.sendmail(sender_email, receiver_email, message.as_string())

    speak_and_log("Email sent successfully.")


# to fetch news using newsapi
def news():
    main_url = "https://newsapi.org/v2/top-headlines?country=in&apiKey=3cd88dd06b794c5b8338fdff94a0f0f2"
    main_page = requests.get(main_url).json()
    articles = main_page["articles"]
    head = []
    day = ["first", "second", "third", "fourth", "fifth"]
    for ar in articles:
        head.append(ar["title"])
    for i in range(len(day)):
        speak_and_log(f"today's {day[i]} news is:{head[i]}")


def show_notification(title, message):
    notification.notify(
        title=title,
        message=message,
        app_name='LUNA Notification',
        timeout=50  # The notification will auto3matically disappear after 50 seconds
    )


@eel.expose
def start():
    wish()
    while True:
        if 1:
            query = takecommand().lower()
            # logic building for task
            if "open notepad" in query or "notepad" in query or "note" in query:
                npath = "C:\\Windows\\notepad.exe"
                os.startfile(npath)

            elif "translate" in query or "translate text" in query or "translation" in query :
                speak_and_log("please speak the text you want to translate")
                user_input = takecommand().lower()
                speak_and_log("please tell me the language name (e.g., 'Spanish')")
                target_language_name = takecommand().lower()

                translated_text = translate_text(user_input, target_language_name)
                speak_and_log(f"Translated text: {translated_text}")

            elif "what's the time" in query or "current time" in query or "time" in query:
                tt = datetime.datetime.now().strftime('%I:%M %p')
                speak_and_log(f"The time is {tt}")

            elif "what is my current location" in query or "where am i" in query or "current city name" in query or "location" in query:
                speak("please wait,let me check")
                currenloc = get_current_location()
                show_notification("current location", f"your location is {currenloc.city} city in {currenloc.country} ")
                log_system_message(f"your location is {currenloc.city} city in {currenloc.country} ")
                speak_and_log(f"your location is {currenloc.city} city in {currenloc.country} ")
            elif "take a screenshot" in query or "screenshot" in query or "capture snapshot" in query or "screen image" in query:
                take_screenshot_and_save()
            elif "opem command prompt" in query or "cmd" in query or "command line" in query:
                os.system("start cmd")
            elif "open camera" in query or "camera" in query or "take a picture" in query or "picture" in query:
                cam = cv2.VideoCapture(0)
                while True:
                    ret, img = cam.read()
                    cv2.imshow('webcam', img)
                    k = cv2.waitKey(50)
                    if k == 27:
                        break
                    speak_and_log("press escape button to close the camera")

                    speak_and_log("do you want to click picture?")
                    res = takecommand().lower()
                    if res == "yes":

                        timestamp = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
                        file_name = f"captured_{timestamp}.png"

                        # Save the screenshot to the desktop
                        desktop_path = "C:\\Users\\Mayuri\\Pictures"
                        file_path = f"{desktop_path}/{file_name}"
                        cv2.imwrite(file_path, img)
                        speak_and_log("picture saved successfully")
                        show_notification("picture saved", "picture saved successfully")
                    else:
                        break
                cam.release()
                cv2.destroyAllWindows()
            elif "open word" in query or "start word" in query or "microsoft word" in query or "word" in query:
                open_word()
            elif "open excel" in query or "start excel" in query or "microsoft excel" in query or "excel" in query:
                open_excel()
            elif "open powerpoint" in query or "start powerpoint" in query or "microsoft powerpoint" in query or "ppt" in query:
                open_powerpoint()

            elif "open vs code" in query or "visual studio code" in query or "code editor" in query:
                open_vscode()
            elif "open pdf file" in query or "pdf file" in query or "pdf" in query:
                # Input the PDF file name
                speak_and_log("please enter the file name you want to open with .pdf extension")
                user_input_file_name = eel.getFilename()()
                # Search for and open the specified PDF file
                pdf_file_path = find_file(user_input_file_name)
                speak_and_log(f"opened the pdf file {user_input_file_name}")
                # Open the PDF file if found
                open_pdf(pdf_file_path)
            elif "open spotify" in query or "spotify" in query or "music player" in query:
                spotify_url = "https://open.spotify.com/"
                webbrowser.open(spotify_url)
            elif "ip address" in query:
                ip = get("https://api.ipify.org").text
                speak_and_log(f"Your IP address is {ip}")
            elif "check internet speed" in query or "internet speed" in query:
                speak_and_log("please wait while i test the internet speed and fetch the data")
                st = speedtest.Speedtest()
                dl = st.download() / 10 ** 6
                up = st.upload() / 10 ** 6
                speak_and_log(f"your internet have {dl:.2f}Mbps downloading speed and {up:.2f}Mbps uploading speed")

            elif "wikipedia" in query or "search wikipedia" in query:
                speak_and_log("What should i search on wikipedia")
                queryy = takecommand().lower()
                results = wikipedia.summary(queryy, sentences=2)
                speak("according to wikipedia...")
                speak_and_log(results)
            elif "open youtube" in query or "youtube" in query:
                webbrowser.open("www.youtube.com")

            elif "play song on youtube" in query or "play song" in query or "start music in youtube" in query:
                speak_and_log("What should I play on youtube")
                play_cont = takecommand().lower()
                kit.playonyt(play_cont)

            elif "play games" in query or "games" in query:
                speak_and_log("Which game do you want to play")
                speak_and_log("1. rock,paper and scissors")
                speak_and_log("2. guess the number")
                speak_and_log("3. riddle")
                game_choice = takecommand().lower()
                if "rock" in game_choice:
                    gamess.play_game()
                elif "guess" in game_choice or "number" in game_choice:
                    gamess.guess_the_number()
                elif "riddle" in game_choice:
                    gamess.get_indian_riddle()
            elif "open google" in query or "google" in query:
                webbrowser.open("www.google.com")
            elif "open calculator" in query or "calculator" in query:
                open_calculator()
            elif "send a whatsapp message" in query or "whatsapp message" in query or "message on whatsapp" in query or "whatsapp" in query:
                speak_and_log("Tell me the contact number to send message with country code:")
                contact = eel.getContact()()
                speak_and_log("what is the content of message")
                msgcont = takecommand().lower()
                kit.sendwhatmsg_instantly(contact, msgcont, 20, 3)
                speak_and_log("Message has been sent successfully!!")

            elif "send email" in query:
                try:
                    speak_and_log("Please enter the recipient email address")
                    recipient_email = eel.getEmail()()

                    # Create the email message
                    speak_and_log("What is the subject of email")
                    subject = takecommand().lower()
                    speak_and_log("What is the body of email")
                    body = takecommand().lower()
                    sendEmail(recipient_email, subject, body)
                    speak_and_log("Email has been sent")
                except Exception as e:
                    speak_and_log(e)
                    speak_and_log("Sorry!I am unable to send email")

            elif "send file via email" in query or "send file using email" in query or "email this file" in query:
                speak_and_log("please enter the recipient's email address")
                receiver_email = eel.getEmail()()
                speak_and_log("please enter the file name with extension")
                file_name = eel.getFilename()()
                send_fileemail(receiver_email, file_name)


            elif "create new word file" in query or "create word document" in query:
                speak_and_log("Tell me the file name for the new Word document (without extension): ")
                user_input = takecommand().lower
                create_and_open_word_document(user_input)

            elif "create new excel file" in query:
                speak_and_log("Enter the file name for the new Excel file (without extension): ")
                user_input = takecommand().lower
                create_and_open_excel(user_input)

            elif "create new powerpoint presentation" in query or "create new ppt" in query:
                speak_and_log("Enter the file name for the new PowerPoint file (without extension): ")
                user_input = takecommand().lower()
                create_and_open_powerpoint(user_input)

            elif "create new python file" in query or "create python file" in query:
                speak_and_log("Enter the file name (without extension) for the new Python file: ")
                user_input = takecommand().lower()
                create_and_open_python_file(user_input)

            elif "tell me a joke" in query or "joke" in query:
                joke = pyjokes.get_joke()
                speak_and_log(joke)

            elif "do some calculation" in query or "calculate this" in query:
                speak_and_log("say what you want to calculate,example: 3 plus 3")
                my_string = takecommand().lower()

                def get_operator_n(op):
                    return {
                        '+': operator.add,
                        '-': operator.sub,
                        'x': operator.mul,
                        'divided': operator.__truediv__,

                    }[op]

                def eval_binary_expr(op1, oper, op2):
                    op1, op2 = int(op1), int(op2)
                    return get_operator_n(oper)(op1, op2)

                speak_and_log("your result is")
                speak_and_log(eval_binary_expr(*(my_string.split())))

            elif "tell me news" in query or "news" in query:
                speak_and_log("Please wait ,fetching the latest news")
                news()



            elif "what's the temperature of city" in query or "what is todays temprature" in query or "temprature" in query:
                location = get_current_location()
                loc = location.city
                search = (f"temperature in {loc}")
                url = f"https://www.google.com//search?q={search}"
                r = requests.get(url)
                data = BeautifulSoup(r.text, "html.parser")
                temp = data.find("div", class_="BNeawe").text
                speak_and_log(f"current temperature of {loc} is {temp}")

            elif "what's the battery percentage" in query or "battery percentage" in query or "battery percent" in query or "charging remaining" in query or "charging percent" in query or "battery status" in query:
                battery = psutil.sensors_battery()
                percentage = battery.percent
                power_plugged = battery.power_plugged
                if power_plugged:
                    status = "charging"
                else:
                    status = "Discharging"
                speak_and_log(f"your system have {percentage} percent battery,and it is currently {status}")
                if percentage >= 75:
                    speak_and_log(" your system have enough power to continue your work")
                elif percentage >= 40:
                    speak_and_log("you should connect your system to charging point to charge the battery")
                elif percentage <= 15:
                    speak_and_log(
                        "your system have very lo power,please connect the charger otherwise the system will shutdown soon")

            elif "switch window" in query or "change window" in query or "change the screen" in query or "change the window" in query:
                pyautogui.keyDown("alt")
                pyautogui.press("tab")
                time.sleep(1)
                pyautogui.keyUp("alt")

            elif "shut down the system" in query:
                os.system("shutdown /s /t 5")

            elif "restart the system" in query or "restart system" in query or "restart the machine" in query or "restart machine" in query:
                os.system("shutdown /r /t 5")

            elif "turn on num lock" in query or "num lock" in query:
                pyautogui.press("numlock")

            elif "turn on caps lock" in query or "caps lock" in query:
                pyautogui.press("capslock")

            elif "volume up" in query or "increase volume" in query or "increase sound" in query:
                pyautogui.press("volumeup")

            elif "volume down" in query or "decrease volume" in query or "decrease sound" in query or "lower down sound" in query:
                pyautogui.press("volumedown")

            elif "volume mute" in query or "mute" in query:
                pyautogui.press("volumemute")

            elif "select all " in query:
                pyautogui.keyDown("ctrl")
                pyautogui.press("a")
                pyautogui.keyUp("ctrl")
                if "copy it" in query or "copy" in query:
                    pyautogui.keyDown("ctrl")
                    pyautogui.press("c")
                    pyautogui.keyUp("ctrl")
                elif "cut it" in query or "cut" in query:
                    pyautogui.keyDown("ctrl")
                    pyautogui.press("x")
                    pyautogui.keyUp("ctrl")
                elif "delete it" in query:
                    pyautogui.press("delete")
                else:
                    pass
            elif "copy" in query:
                pyautogui.keyDown("ctrl")
                pyautogui.press("c")
                pyautogui.keyUp("ctrl")
            elif "cut" in query:
                pyautogui.keyDown("ctrl")
                pyautogui.press("x")
                pyautogui.keyUp("ctrl")
            elif "paste" in query:
                pyautogui.keyDown("ctrl")
                pyautogui.press("v")
                pyautogui.keyUp("ctrl")
            elif "hello" in query:
                speak_and_log("hello,may i help you with something")

            elif "how are you" in query:
                speak_and_log("i am fine , what about you")

            elif "also good" in query:
                speak_and_log("that's great to hear from you")

            elif "you can sleep" in query:
                speak_and_log("okay,i am going to sleep you can call me anytime")
                sys.exit()

            elif "sleep the system" in query:
                os.system("shutdown /h")

            elif "no thanks" in query:
                speak_and_log("Thanks for using me")
                return "exit"
            speak_and_log("Do you have any other tasks?")
            response = takecommand().lower()
            if "no" in response:
                speak_and_log("Thanks for using me")
                return "exit"
            elif "yes" in response:
                speak_and_log("please tell me the task")
                additional_tasks = False
            else:
                speak_and_log("I didn't understand that. Please say 'yes' or 'no'.")


if __name__ == "__main__":
    # Start the Eel application
    eel.start('templates\\index.html', size=(400, 300), mode='chrome')
